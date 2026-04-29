"""
Helper library for applying Winningtemp brand palette and typography to python-pptx
shapes/runs. Import from your own scripts; this file is a library, not a CLI.

    from apply_brand_palette import BrandColor, BrandFont, brand_run, brand_fill, brand_card

    title_run.font.name = BrandFont.HEADING
    brand_fill(shape, BrandColor.PURPLE)
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.util import Pt

SKILL_DIR = Path(__file__).resolve().parent.parent
ASSETS = SKILL_DIR / "assets"
LOGOS = ASSETS / "logos"
ICONS = ASSETS / "icons" / "black"
TEMPLATES = ASSETS / "templates"


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


with (ASSETS / "color-palette.json").open() as f:
    _palette = json.load(f)
with (ASSETS / "typography.json").open() as f:
    _typography = json.load(f)


class BrandColor:
    """Brand palette as RGBColor objects, ready for python-pptx."""
    PURPLE       = _hex_to_rgb(_palette["primary"]["purple"]["hex"])         # #7F2696
    LIGHT_PURPLE = _hex_to_rgb(_palette["primary"]["light_purple"]["hex"])   # #DDD3F2
    GRAY         = _hex_to_rgb(_palette["primary"]["gray"]["hex"])           # #F1F1F1
    PALE_PURPLE  = _hex_to_rgb(_palette["secondary"]["extended_purple_light"]["hex"])  # #F2E9F3
    WAVE         = _hex_to_rgb(_palette["secondary"]["wave"]["hex"])         # #C0DAD8
    WARNING_RED  = _hex_to_rgb(_palette["secondary"]["red_warning"]["hex"])  # #D64545
    BLACK        = _hex_to_rgb(_palette["support"]["black"]["hex"])          # #18181B
    WHITE        = _hex_to_rgb(_palette["support"]["white"]["hex"])          # #FFFFFF


class BrandFont:
    HEADING = _typography["headings"]["family"]    # "Poppins"
    BODY    = _typography["body"]["family"]        # "Inter"


@dataclass(frozen=True)
class TypeStyle:
    name: str
    size_pt: int
    bold: bool = False
    color: RGBColor = BrandColor.BLACK


# Common pre-built type styles
H1            = TypeStyle(BrandFont.HEADING, 48, bold=True,  color=BrandColor.BLACK)
H1_ON_PURPLE  = TypeStyle(BrandFont.HEADING, 48, bold=True,  color=BrandColor.WHITE)
H2            = TypeStyle(BrandFont.HEADING, 24, bold=False, color=BrandColor.BLACK)
H2_ON_PURPLE  = TypeStyle(BrandFont.HEADING, 24, bold=False, color=BrandColor.WHITE)
BODY          = TypeStyle(BrandFont.BODY,    14, bold=False, color=BrandColor.BLACK)
BODY_ON_PURPLE= TypeStyle(BrandFont.BODY,    14, bold=False, color=BrandColor.WHITE)
CAPTION       = TypeStyle(BrandFont.BODY,    10, bold=False, color=BrandColor.BLACK)


def brand_run(run, style: TypeStyle) -> None:
    """Apply a TypeStyle to a python-pptx text Run in place."""
    run.font.name = style.name
    run.font.size = Pt(style.size_pt)
    run.font.bold = style.bold
    run.font.color.rgb = style.color


def brand_paragraph(paragraph, text: str, style: TypeStyle) -> None:
    """Set a paragraph's text and apply a brand style."""
    paragraph.text = text
    if paragraph.runs:
        brand_run(paragraph.runs[0], style)


def brand_fill(shape, color: RGBColor) -> None:
    """Fill a shape with a brand color and remove its border."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def brand_card(slide, left, top, width, height,
               fill_color: RGBColor = BrandColor.WHITE,
               accent_color: RGBColor | None = BrandColor.PURPLE,
               accent_height_pt: int = 4):
    """
    Add a card-style container: rounded fill + optional top accent bar.
    Uses MSO_SHAPE.ROUNDED_RECTANGLE (1) for the body, and MSO_SHAPE.RECTANGLE (1) for the bar.
    """
    from pptx.enum.shapes import MSO_SHAPE
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    brand_fill(card, fill_color)
    if accent_color is not None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(accent_height_pt))
        brand_fill(bar, accent_color)
    return card


def add_logo(slide, left, top, width, color: str = "deep-purple", with_tagline: bool = False):
    """
    Add the Winningtemp logo to a slide.

    color: "deep-purple" (default) | "white" | "black"
    with_tagline: True to use the variant with the SUCCEED TOGETHER tagline
    """
    suffix = " tagline" if with_tagline else " (Without tagline)"
    color_folder_map = {
        "deep-purple": "Deep purple",
        "white":       "White",
        "black":       "Black",
    }
    if color not in color_folder_map:
        raise ValueError(f"color must be one of {list(color_folder_map)}; got {color!r}")
    folder = color_folder_map[color]

    if with_tagline:
        path = LOGOS / "variations" / "Winningtemp Logo" / folder / f"Winningtemp_RGB {folder}.png"
        if not path.exists():
            # Fallback to the explicit (tagline) folder
            path = LOGOS / "variations" / "Winningtemp Logo (tagline)" / f"{folder} (tagline)" / f"Winningtemp_RGB {folder} tagline.png"
    else:
        path = LOGOS / "variations" / "Winningtemp Logo (Without tagline)" / f"{folder} (Without tagline)" / f"Winningtemp_RGB {folder} (Without tagline).png"
        if not path.exists():
            path = LOGOS / "standard" / f"Winningtemp_RGB {folder} (Without tagline).png"

    if not path.exists():
        raise FileNotFoundError(f"Logo not found: {path}")
    return slide.shapes.add_picture(str(path), left, top, width=width)


def add_icon(slide, icon_name: str, left, top, width):
    """
    Add an official Winningtemp icon. icon_name is the file's display name,
    case- and extension-insensitive (e.g. "insight" → "Insight.png").
    Looks up from assets/icons/black/.
    """
    target = icon_name.lower().removesuffix(".png").strip()
    matches = [p for p in ICONS.glob("*.png") if p.stem.lower() == target]
    if not matches:
        # Loose contains-match fallback
        matches = [p for p in ICONS.glob("*.png") if target in p.stem.lower()]
    if not matches:
        available = sorted(p.stem for p in ICONS.glob("*.png"))
        raise FileNotFoundError(
            f"No icon named {icon_name!r} in {ICONS}. "
            f"Try one of: {', '.join(available[:20])}..."
        )
    return slide.shapes.add_picture(str(matches[0]), left, top, width=width)


def open_branded_template():
    """
    Open the official Winningtemp PowerPoint template as a Presentation.
    Use this as the base for any new deck instead of starting from scratch.
    """
    from pptx import Presentation
    template = TEMPLATES / "winningtemp-presentation-template.pptx"
    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")
    return Presentation(str(template))
