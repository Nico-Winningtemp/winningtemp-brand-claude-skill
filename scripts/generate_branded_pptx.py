#!/usr/bin/env python3
"""
Reusable branded PPTX builder. Replaces ad-hoc deck generators (like the
generic-themed generate_slides.py from March 2026) with one that produces
on-brand decks by default.

Two ways to use:

1) CLI quick-test:
       python3 scripts/generate_branded_pptx.py --title "My Deck" --output my-deck.pptx

2) Library — import build_deck() and pass a list of slide specs:

       from generate_branded_pptx import build_deck, TitleSlide, KPISlide, ContentSlide

       build_deck("report.pptx", slides=[
           TitleSlide("Q1 Results", "April 2026 | Winningtemp"),
           KPISlide("Engagement", [("MAU", "12,456"), ("CSAT", "94%")]),
           ContentSlide("Findings", ["Point one.", "Point two."]),
       ])

The script always opens the official template
(assets/templates/winningtemp-presentation-template.pptx) as the base, so
master colors, fonts, and slide masters are correct.
"""

from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass, field
from pathlib import Path

# Make sibling helper importable
sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from apply_brand_palette import (
    BrandColor, BrandFont,
    H1, H1_ON_PURPLE, H2, H2_ON_PURPLE, BODY, BODY_ON_PURPLE, CAPTION,
    brand_run, brand_fill, brand_card, add_logo, add_icon,
    open_branded_template,
)


# ---------- slide spec dataclasses ----------

@dataclass
class TitleSlide:
    title: str
    subtitle: str = ""

@dataclass
class SectionSlide:
    label: str               # e.g. "01. Findings"
    title: str               # e.g. "What we learned"

@dataclass
class KPISlide:
    title: str
    kpis: list[tuple[str, str]]   # list of (label, value)
    note: str = ""

@dataclass
class ContentSlide:
    title: str
    bullets: list[str]
    body: str = ""

@dataclass
class TwoColumnSlide:
    title: str
    left_title: str
    left_body: str
    right_title: str
    right_body: str

@dataclass
class EndSlide:
    headline: str = "Thank you"
    subline: str = ""


SlideSpec = TitleSlide | SectionSlide | KPISlide | ContentSlide | TwoColumnSlide | EndSlide


# ---------- text helpers ----------

def _add_text(slide, left, top, width, height, text, style, alignment=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    if p.runs:
        brand_run(p.runs[0], style)
    return box


def _add_bullets(slide, left, top, width, height, bullets, style):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.space_after = Pt(8)
        if p.runs:
            brand_run(p.runs[0], style)
    return box


# ---------- slide renderers ----------

def _render_title(prs, spec: TitleSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    # Purple background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    brand_fill(bg, BrandColor.PURPLE)

    _add_text(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.6),
              spec.title, H1_ON_PURPLE, alignment=PP_ALIGN.LEFT)
    if spec.subtitle:
        _add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
                  spec.subtitle, BODY_ON_PURPLE)

    # White logo bottom-right
    add_logo(slide, prs.slide_width - Inches(2.2), prs.slide_height - Inches(0.9),
             Inches(1.6), color="white")


def _render_section(prs, spec: SectionSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    brand_fill(bg, BrandColor.LIGHT_PURPLE)

    _add_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.7),
              spec.label, H2)
    _add_text(slide, Inches(1), Inches(3.4), Inches(11), Inches(1.5),
              spec.title, H1)


def _render_kpi(prs, spec: KPISlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
              spec.title, H1)

    n = len(spec.kpis)
    if n == 0:
        return
    gap = Inches(0.25)
    total_w = prs.slide_width - Inches(1.2) - gap * (n - 1)
    card_w = total_w / n
    card_h = Inches(2.0)
    y = Inches(2.0)

    for i, (label, value) in enumerate(spec.kpis):
        x = Inches(0.6) + (card_w + gap) * i
        brand_card(slide, x, y, card_w, card_h,
                   fill_color=BrandColor.GRAY,
                   accent_color=BrandColor.PURPLE)
        _add_text(slide, x + Inches(0.2), y + Inches(0.4),
                  card_w - Inches(0.4), Inches(0.8),
                  value, H1, alignment=PP_ALIGN.CENTER)
        _add_text(slide, x + Inches(0.2), y + Inches(1.3),
                  card_w - Inches(0.4), Inches(0.5),
                  label, CAPTION, alignment=PP_ALIGN.CENTER)

    if spec.note:
        _add_text(slide, Inches(0.6), Inches(5.0), Inches(11), Inches(0.6),
                  spec.note, CAPTION)


def _render_content(prs, spec: ContentSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
              spec.title, H1)
    if spec.body:
        _add_text(slide, Inches(0.6), Inches(1.4), Inches(11), Inches(0.8),
                  spec.body, BODY)
    if spec.bullets:
        top_y = Inches(2.4) if spec.body else Inches(1.5)
        _add_bullets(slide, Inches(0.6), top_y, Inches(11),
                     prs.slide_height - top_y - Inches(0.5),
                     spec.bullets, BODY)


def _render_two_column(prs, spec: TwoColumnSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
              spec.title, H1)

    col_w = (prs.slide_width - Inches(1.5)) / 2
    col_h = prs.slide_height - Inches(2.0)

    brand_card(slide, Inches(0.6), Inches(1.6), col_w, col_h,
               fill_color=BrandColor.GRAY, accent_color=BrandColor.PURPLE)
    _add_text(slide, Inches(0.9), Inches(1.9), col_w - Inches(0.6), Inches(0.6),
              spec.left_title, H2)
    _add_text(slide, Inches(0.9), Inches(2.6), col_w - Inches(0.6), col_h - Inches(1.2),
              spec.left_body, BODY)

    right_x = Inches(0.6) + col_w + Inches(0.3)
    brand_card(slide, right_x, Inches(1.6), col_w, col_h,
               fill_color=BrandColor.LIGHT_PURPLE, accent_color=BrandColor.PURPLE)
    _add_text(slide, right_x + Inches(0.3), Inches(1.9), col_w - Inches(0.6), Inches(0.6),
              spec.right_title, H2)
    _add_text(slide, right_x + Inches(0.3), Inches(2.6), col_w - Inches(0.6), col_h - Inches(1.2),
              spec.right_body, BODY)


def _render_end(prs, spec: EndSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    brand_fill(bg, BrandColor.PURPLE)
    _add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.4),
              spec.headline, H1_ON_PURPLE, alignment=PP_ALIGN.CENTER)
    if spec.subline:
        _add_text(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
                  spec.subline, BODY_ON_PURPLE, alignment=PP_ALIGN.CENTER)


_DISPATCH = {
    TitleSlide: _render_title,
    SectionSlide: _render_section,
    KPISlide: _render_kpi,
    ContentSlide: _render_content,
    TwoColumnSlide: _render_two_column,
    EndSlide: _render_end,
}


# ---------- public API ----------

def _wipe_template_slides(prs) -> None:
    """Remove all existing slides from the template presentation cleanly.
    Removes both the sldId entries AND the underlying slide parts/rels so the
    output .pptx doesn't contain duplicates.
    """
    sld_id_lst = prs.slides._sldIdLst
    rels = prs.part.rels
    for sld_id in list(sld_id_lst):
        rId = sld_id.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
        slide_part = rels[rId].target_part
        prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)
        # Also drop the slide part from the package so its XML doesn't get re-emitted
        try:
            del prs.part.package._parts_by_partname[slide_part.partname]
        except (AttributeError, KeyError):
            pass


def build_deck(output_path: str | Path, slides: list[SlideSpec]) -> Path:
    """Build a branded deck from a list of slide specs and save to output_path."""
    prs = open_branded_template()
    _wipe_template_slides(prs)
    for spec in slides:
        renderer = _DISPATCH.get(type(spec))
        if renderer is None:
            raise TypeError(f"Unknown slide spec type: {type(spec).__name__}")
        renderer(prs, spec)
    output_path = Path(output_path).resolve()
    prs.save(str(output_path))
    return output_path


# ---------- CLI ----------

def _demo_slides(title: str) -> list[SlideSpec]:
    return [
        TitleSlide(title=title, subtitle="Generated by winningtemp-brand skill"),
        SectionSlide(label="01. Overview", title="What's inside"),
        KPISlide("Headline KPIs",
                 [("Engagement", "82%"), ("eNPS", "+34"), ("Response", "76%")],
                 note="Sample data — replace before sharing."),
        ContentSlide("Findings",
                     bullets=[
                         "Engagement is highest in teams with weekly check-ins.",
                         "eNPS responses are warmest when managers reply within 48 hours.",
                         "Frontline teams need a different channel — mobile-first.",
                     ],
                     body="Three patterns from this quarter that change how we think about manager support."),
        TwoColumnSlide("What worked / What didn't",
                       left_title="Worked", left_body="Weekly pulse cadence. Short, focused questions.",
                       right_title="Didn't", right_body="Quarterly deep-dives. Too long between signals."),
        EndSlide(headline="Questions?", subline="hello@winningtemp.com"),
    ]


def main():
    ap = argparse.ArgumentParser(description="Build a branded Winningtemp PPTX.")
    ap.add_argument("--title", default="Winningtemp Demo Deck", help="Title for the demo deck")
    ap.add_argument("--output", default="demo-branded-deck.pptx", help="Output .pptx path")
    args = ap.parse_args()

    out = build_deck(args.output, _demo_slides(args.title))
    print(f"✓ Built deck: {out}")
    print(f"  Slides: 6 (title, section, KPI, content, two-column, end)")
    print(f"  Open with: open '{out}'")


if __name__ == "__main__":
    main()
